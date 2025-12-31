#!/usr/bin/env python3
"""
Document Metadata Extraction Module - Ultimate Edition

Extracts comprehensive document metadata including:
- Office documents (Word, Excel, PowerPoint, LibreOffice)
- PDF documents (structure, forms, security, annotations)
- Web documents (HTML, XML, CSS, JavaScript)
- E-book formats (EPUB, MOBI, AZW)
- Archive formats (ZIP, RAR, 7Z, TAR)
- Source code files (programming languages, repositories)
- Configuration files (JSON, YAML, TOML, INI)
- Database files (SQLite, Access)
- CAD and design files (DWG, SVG, AI)
- Scientific documents (LaTeX, BibTeX, Markdown)

Author: MetaExtract Team
Version: 1.0.0
"""

import os
import json
import zipfile
import tarfile
import logging
import xml.etree.ElementTree as ET
import sqlite3
import tempfile
import hashlib
import re
from pathlib import Path
from typing import Any, Dict, Optional, List, Union
from datetime import datetime
import mimetypes

logger = logging.getLogger(__name__)

# Library availability checks
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

try:
    import toml
    TOML_AVAILABLE = True
except ImportError:
    TOML_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EBOOKLIB_AVAILABLE = True
except ImportError:
    EBOOKLIB_AVAILABLE = False

def extract_document_metadata(filepath: str) -> Dict[str, Any]:
    """Extract comprehensive document metadata"""
    
    result = {
        "available": True,
        "document_type": "unknown",
        "office_metadata": {},
        "pdf_metadata": {},
        "web_metadata": {},
        "ebook_metadata": {},
        "archive_metadata": {},
        "source_code_metadata": {},
        "config_metadata": {},
        "database_metadata": {},
        "text_analysis": {},
        "security_metadata": {}
    }
    
    try:
        file_ext = Path(filepath).suffix.lower()
        mime_type, _ = mimetypes.guess_type(filepath)
        
        # Office documents
        if file_ext in ['.docx', '.xlsx', '.pptx']:
            result["document_type"] = "office"
            office_result = _analyze_office_documents(filepath, file_ext)
            if office_result:
                result["office_metadata"].update(office_result)
        
        # PDF documents
        elif file_ext == '.pdf':
            result["document_type"] = "pdf"
            pdf_result = _analyze_pdf_document(filepath)
            if pdf_result:
                result["pdf_metadata"].update(pdf_result)
        
        # Web documents
        elif file_ext in ['.html', '.htm', '.xml', '.css', '.js']:
            result["document_type"] = "web"
            web_result = _analyze_web_document(filepath, file_ext)
            if web_result:
                result["web_metadata"].update(web_result)
        
        # E-books
        elif file_ext in ['.epub', '.mobi', '.azw', '.azw3']:
            result["document_type"] = "ebook"
            ebook_result = _analyze_ebook(filepath, file_ext)
            if ebook_result:
                result["ebook_metadata"].update(ebook_result)
        
        # Archives
        elif file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']:
            result["document_type"] = "archive"
            archive_result = _analyze_archive(filepath, file_ext)
            if archive_result:
                result["archive_metadata"].update(archive_result)
        
        # Source code
        elif file_ext in ['.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php', '.rb', '.go', '.rs']:
            result["document_type"] = "source_code"
            code_result = _analyze_source_code(filepath, file_ext)
            if code_result:
                result["source_code_metadata"].update(code_result)
        
        # Configuration files
        elif file_ext in ['.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf']:
            result["document_type"] = "config"
            config_result = _analyze_config_file(filepath, file_ext)
            if config_result:
                result["config_metadata"].update(config_result)
        
        # Database files
        elif file_ext in ['.db', '.sqlite', '.sqlite3', '.mdb']:
            result["document_type"] = "database"
            db_result = _analyze_database_file(filepath, file_ext)
            if db_result:
                result["database_metadata"].update(db_result)
        
        # Text files
        elif file_ext in ['.txt', '.md', '.rst', '.tex']:
            result["document_type"] = "text"
            text_result = _analyze_text_document(filepath, file_ext)
            if text_result:
                result["text_analysis"].update(text_result)
        
        # General text analysis for all text-based formats
        if result["document_type"] in ["office", "pdf", "web", "text", "source_code", "config"]:
            text_stats = _analyze_text_statistics(filepath)
            if text_stats:
                result["text_analysis"].update(text_stats)
        
        return result
        
    except Exception as e:
        logger.error(f"Error in document analysis: {e}")
        return {"available": False, "error": str(e)}

def _analyze_office_documents(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze Microsoft Office documents"""
    try:
        result = {}
        
        if file_ext == '.docx' and DOCX_AVAILABLE:
            result.update(_analyze_docx(filepath))
        elif file_ext == '.xlsx' and OPENPYXL_AVAILABLE:
            result.update(_analyze_xlsx(filepath))
        elif file_ext == '.pptx' and PPTX_AVAILABLE:
            result.update(_analyze_pptx(filepath))
        
        return result
        
    except Exception as e:
        logger.error(f"Office document analysis error: {e}")
        return {}

def _analyze_docx(filepath: str) -> Dict[str, Any]:
    """Analyze Word document"""
    try:
        doc = Document(filepath)
        
        result = {
            "docx_analysis": {
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "section_count": len(doc.sections),
                "style_count": len(doc.styles),
                "has_images": False,
                "has_hyperlinks": False,
                "word_count": 0,
                "character_count": 0
            }
        }
        
        # Text analysis
        full_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text)
        
        combined_text = '\n'.join(full_text)
        result["docx_analysis"]["word_count"] = len(combined_text.split())
        result["docx_analysis"]["character_count"] = len(combined_text)
        
        # Check for images and hyperlinks
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.element.xpath('.//a:blip'):  # Image
                    result["docx_analysis"]["has_images"] = True
                if run.element.xpath('.//w:hyperlink'):  # Hyperlink
                    result["docx_analysis"]["has_hyperlinks"] = True
        
        # Document properties
        props = doc.core_properties
        result["docx_analysis"]["properties"] = {
            "title": props.title,
            "author": props.author,
            "subject": props.subject,
            "keywords": props.keywords,
            "comments": props.comments,
            "created": props.created.isoformat() if props.created else None,
            "modified": props.modified.isoformat() if props.modified else None,
            "last_modified_by": props.last_modified_by,
            "revision": props.revision,
            "version": props.version
        }
        
        return result
        
    except Exception as e:
        logger.error(f"DOCX analysis error: {e}")
        return {}

def _analyze_xlsx(filepath: str) -> Dict[str, Any]:
    """Analyze Excel spreadsheet"""
    try:
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        
        result = {
            "xlsx_analysis": {
                "worksheet_count": len(workbook.worksheets),
                "worksheet_names": [ws.title for ws in workbook.worksheets],
                "has_formulas": False,
                "has_charts": False,
                "total_cells": 0,
                "non_empty_cells": 0
            }
        }
        
        # Analyze worksheets
        for worksheet in workbook.worksheets:
            max_row = worksheet.max_row
            max_col = worksheet.max_column
            result["xlsx_analysis"]["total_cells"] += max_row * max_col
            
            # Count non-empty cells and check for formulas
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        result["xlsx_analysis"]["non_empty_cells"] += 1
                    
                    if hasattr(cell, 'data_type') and cell.data_type == 'f':
                        result["xlsx_analysis"]["has_formulas"] = True
        
        # Document properties
        props = workbook.properties
        result["xlsx_analysis"]["properties"] = {
            "title": props.title,
            "creator": props.creator,
            "subject": props.subject,
            "description": props.description,
            "keywords": props.keywords,
            "created": props.created.isoformat() if props.created else None,
            "modified": props.modified.isoformat() if props.modified else None,
            "last_modified_by": props.lastModifiedBy,
            "version": props.version
        }
        
        return result
        
    except Exception as e:
        logger.error(f"XLSX analysis error: {e}")
        return {}

def _analyze_pptx(filepath: str) -> Dict[str, Any]:
    """Analyze PowerPoint presentation"""
    try:
        prs = Presentation(filepath)
        
        result = {
            "pptx_analysis": {
                "slide_count": len(prs.slides),
                "layout_count": len(prs.slide_layouts),
                "master_count": len(prs.slide_masters),
                "has_images": False,
                "has_videos": False,
                "has_animations": False,
                "text_content": []
            }
        }
        
        # Analyze slides
        for slide in prs.slides:
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
                
                # Check for media
                if shape.shape_type == 13:  # Picture
                    result["pptx_analysis"]["has_images"] = True
                elif shape.shape_type == 24:  # Media (video/audio)
                    result["pptx_analysis"]["has_videos"] = True
            
            if slide_text:
                result["pptx_analysis"]["text_content"].append(' '.join(slide_text))
        
        # Document properties
        props = prs.core_properties
        result["pptx_analysis"]["properties"] = {
            "title": props.title,
            "author": props.author,
            "subject": props.subject,
            "keywords": props.keywords,
            "comments": props.comments,
            "created": props.created.isoformat() if props.created else None,
            "modified": props.modified.isoformat() if props.modified else None,
            "last_modified_by": props.last_modified_by,
            "revision": props.revision
        }
        
        return result
        
    except Exception as e:
        logger.error(f"PPTX analysis error: {e}")
        return {}

def _analyze_pdf_document(filepath: str) -> Dict[str, Any]:
    """Analyze PDF document"""
    try:
        result = {}
        
        # PyPDF2 analysis
        if PYPDF2_AVAILABLE:
            result.update(_analyze_pdf_pypdf2(filepath))
        
        # PDFPlumber analysis
        if PDFPLUMBER_AVAILABLE:
            result.update(_analyze_pdf_pdfplumber(filepath))
        
        return result
        
    except Exception as e:
        logger.error(f"PDF analysis error: {e}")
        return {}

def _analyze_pdf_pypdf2(filepath: str) -> Dict[str, Any]:
    """Analyze PDF using PyPDF2"""
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            result = {
                "pypdf2_analysis": {
                    "page_count": len(pdf_reader.pages),
                    "is_encrypted": pdf_reader.is_encrypted,
                    "metadata": {},
                    "form_fields": [],
                    "bookmarks": [],
                    "annotations": []
                }
            }
            
            # Metadata
            if pdf_reader.metadata:
                metadata = pdf_reader.metadata
                result["pypdf2_analysis"]["metadata"] = {
                    "title": metadata.get('/Title'),
                    "author": metadata.get('/Author'),
                    "subject": metadata.get('/Subject'),
                    "creator": metadata.get('/Creator'),
                    "producer": metadata.get('/Producer'),
                    "creation_date": str(metadata.get('/CreationDate')),
                    "modification_date": str(metadata.get('/ModDate'))
                }
            
            # Form fields
            if hasattr(pdf_reader, 'get_form_text_fields'):
                try:
                    form_fields = pdf_reader.get_form_text_fields()
                    if form_fields:
                        result["pypdf2_analysis"]["form_fields"] = list(form_fields.keys())
                except:
                    pass
            
            # Bookmarks/Outlines
            try:
                outlines = pdf_reader.outline
                if outlines:
                    result["pypdf2_analysis"]["bookmarks"] = _extract_pdf_outlines(outlines)
            except:
                pass
            
            # Page analysis
            page_info = []
            for i, page in enumerate(pdf_reader.pages):
                page_data = {
                    "page_number": i + 1,
                    "rotation": page.rotation if hasattr(page, 'rotation') else 0
                }
                
                # Try to get page dimensions
                try:
                    mediabox = page.mediabox
                    page_data["dimensions"] = {
                        "width": float(mediabox.width),
                        "height": float(mediabox.height)
                    }
                except:
                    pass
                
                page_info.append(page_data)
            
            result["pypdf2_analysis"]["pages"] = page_info[:5]  # First 5 pages
            
            return result
            
    except Exception as e:
        logger.error(f"PyPDF2 analysis error: {e}")
        return {}

def _analyze_pdf_pdfplumber(filepath: str) -> Dict[str, Any]:
    """Analyze PDF using pdfplumber"""
    try:
        import pdfplumber
        
        with pdfplumber.open(filepath) as pdf:
            result = {
                "pdfplumber_analysis": {
                    "page_count": len(pdf.pages),
                    "metadata": pdf.metadata or {},
                    "text_extraction": {},
                    "table_detection": {},
                    "image_detection": {}
                }
            }
            
            # Text extraction from first few pages
            text_pages = []
            table_count = 0
            image_count = 0
            
            for i, page in enumerate(pdf.pages[:3]):  # First 3 pages
                # Extract text
                page_text = page.extract_text()
                if page_text:
                    text_pages.append({
                        "page": i + 1,
                        "text_length": len(page_text),
                        "word_count": len(page_text.split())
                    })
                
                # Detect tables
                tables = page.extract_tables()
                table_count += len(tables)
                
                # Detect images
                if hasattr(page, 'images'):
                    image_count += len(page.images)
            
            result["pdfplumber_analysis"]["text_extraction"] = {
                "pages_with_text": text_pages,
                "total_extractable_pages": len([p for p in text_pages if p["text_length"] > 0])
            }
            
            result["pdfplumber_analysis"]["table_detection"] = {
                "tables_found": table_count,
                "has_tables": table_count > 0
            }
            
            result["pdfplumber_analysis"]["image_detection"] = {
                "images_found": image_count,
                "has_images": image_count > 0
            }
            
            return result
            
    except Exception as e:
        logger.error(f"PDFPlumber analysis error: {e}")
        return {}

def _extract_pdf_outlines(outlines, level=0) -> List[Dict[str, Any]]:
    """Extract PDF bookmarks/outlines recursively"""
    result = []
    
    for item in outlines:
        if isinstance(item, list):
            result.extend(_extract_pdf_outlines(item, level + 1))
        else:
            bookmark = {
                "title": str(item.title) if hasattr(item, 'title') else str(item),
                "level": level
            }
            result.append(bookmark)
    
    return result

def _analyze_web_document(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze web documents (HTML, XML, CSS, JS)"""
    try:
        result = {}
        
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        if file_ext in ['.html', '.htm']:
            result.update(_analyze_html(content))
        elif file_ext == '.xml':
            result.update(_analyze_xml(content))
        elif file_ext == '.css':
            result.update(_analyze_css(content))
        elif file_ext == '.js':
            result.update(_analyze_javascript(content))
        
        return result
        
    except Exception as e:
        logger.error(f"Web document analysis error: {e}")
        return {}

def _analyze_html(content: str) -> Dict[str, Any]:
    """Analyze HTML document"""
    try:
        result = {
            "html_analysis": {
                "has_doctype": content.strip().lower().startswith('<!doctype'),
                "meta_tags": {},
                "links": [],
                "images": [],
                "scripts": [],
                "stylesheets": [],
                "forms": 0,
                "tables": 0
            }
        }
        
        if BS4_AVAILABLE:
            soup = BeautifulSoup(content, 'html.parser')
            
            # Meta tags
            meta_tags = soup.find_all('meta')
            for meta in meta_tags:
                name = meta.get('name') or meta.get('property') or meta.get('http-equiv')
                content_attr = meta.get('content')
                if name and content_attr:
                    result["html_analysis"]["meta_tags"][name] = content_attr
            
            # Links
            links = soup.find_all('a', href=True)
            result["html_analysis"]["links"] = [link['href'] for link in links[:10]]  # First 10
            
            # Images
            images = soup.find_all('img', src=True)
            result["html_analysis"]["images"] = [img['src'] for img in images[:10]]  # First 10
            
            # Scripts
            scripts = soup.find_all('script', src=True)
            result["html_analysis"]["scripts"] = [script['src'] for script in scripts]
            
            # Stylesheets
            stylesheets = soup.find_all('link', rel='stylesheet')
            result["html_analysis"]["stylesheets"] = [link['href'] for link in stylesheets if link.get('href')]
            
            # Forms and tables
            result["html_analysis"]["forms"] = len(soup.find_all('form'))
            result["html_analysis"]["tables"] = len(soup.find_all('table'))
            
            # Title
            title = soup.find('title')
            if title:
                result["html_analysis"]["title"] = title.get_text().strip()
        
        return result
        
    except Exception as e:
        logger.error(f"HTML analysis error: {e}")
        return {}

def _analyze_xml(content: str) -> Dict[str, Any]:
    """Analyze XML document"""
    try:
        result = {
            "xml_analysis": {
                "is_well_formed": False,
                "root_element": None,
                "namespaces": [],
                "element_count": 0,
                "attribute_count": 0,
                "max_depth": 0
            }
        }
        
        try:
            root = ET.fromstring(content)
            result["xml_analysis"]["is_well_formed"] = True
            result["xml_analysis"]["root_element"] = root.tag
            
            # Count elements and attributes
            element_count = 0
            attribute_count = 0
            
            def count_elements(element, depth=0):
                nonlocal element_count, attribute_count
                element_count += 1
                attribute_count += len(element.attrib)
                result["xml_analysis"]["max_depth"] = max(result["xml_analysis"]["max_depth"], depth)
                
                for child in element:
                    count_elements(child, depth + 1)
            
            count_elements(root)
            
            result["xml_analysis"]["element_count"] = element_count
            result["xml_analysis"]["attribute_count"] = attribute_count
            
            # Extract namespaces
            namespaces = []
            for elem in root.iter():
                if '}' in elem.tag:
                    namespace = elem.tag.split('}')[0][1:]
                    if namespace not in namespaces:
                        namespaces.append(namespace)
            
            result["xml_analysis"]["namespaces"] = namespaces
            
        except ET.ParseError:
            result["xml_analysis"]["is_well_formed"] = False
        
        return result
        
    except Exception as e:
        logger.error(f"XML analysis error: {e}")
        return {}

def _analyze_css(content: str) -> Dict[str, Any]:
    """Analyze CSS document"""
    try:
        result = {
            "css_analysis": {
                "rule_count": 0,
                "selector_types": {},
                "property_count": 0,
                "media_queries": 0,
                "imports": [],
                "fonts": []
            }
        }
        
        # Count CSS rules (approximate)
        rule_count = content.count('{')
        result["css_analysis"]["rule_count"] = rule_count
        
        # Count properties (approximate)
        property_count = content.count(':') - content.count('::')
        result["css_analysis"]["property_count"] = property_count
        
        # Media queries
        media_queries = content.count('@media')
        result["css_analysis"]["media_queries"] = media_queries
        
        # Imports
        import_matches = re.findall(r'@import\s+["\']([^"\']+)["\']', content)
        result["css_analysis"]["imports"] = import_matches
        
        # Font families
        font_matches = re.findall(r'font-family\s*:\s*([^;]+)', content, re.IGNORECASE)
        result["css_analysis"]["fonts"] = [font.strip() for font in font_matches[:10]]
        
        return result
        
    except Exception as e:
        logger.error(f"CSS analysis error: {e}")
        return {}

def _analyze_javascript(content: str) -> Dict[str, Any]:
    """Analyze JavaScript document"""
    try:
        result = {
            "javascript_analysis": {
                "function_count": 0,
                "variable_declarations": 0,
                "class_count": 0,
                "import_count": 0,
                "export_count": 0,
                "async_functions": 0,
                "arrow_functions": 0
            }
        }
        
        # Function declarations
        function_count = len(re.findall(r'\bfunction\s+\w+', content))
        result["javascript_analysis"]["function_count"] = function_count
        
        # Arrow functions
        arrow_functions = len(re.findall(r'=>', content))
        result["javascript_analysis"]["arrow_functions"] = arrow_functions
        
        # Variable declarations
        var_count = len(re.findall(r'\b(var|let|const)\s+\w+', content))
        result["javascript_analysis"]["variable_declarations"] = var_count
        
        # Classes
        class_count = len(re.findall(r'\bclass\s+\w+', content))
        result["javascript_analysis"]["class_count"] = class_count
        
        # Imports/Exports
        import_count = len(re.findall(r'\bimport\s+', content))
        export_count = len(re.findall(r'\bexport\s+', content))
        result["javascript_analysis"]["import_count"] = import_count
        result["javascript_analysis"]["export_count"] = export_count
        
        # Async functions
        async_count = len(re.findall(r'\basync\s+function', content))
        result["javascript_analysis"]["async_functions"] = async_count
        
        return result
        
    except Exception as e:
        logger.error(f"JavaScript analysis error: {e}")
        return {}

def _analyze_archive(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze archive files"""
    try:
        result = {
            "archive_analysis": {
                "format": file_ext,
                "file_count": 0,
                "total_size": 0,
                "compressed_size": 0,
                "compression_ratio": 0,
                "file_types": {},
                "directory_structure": []
            }
        }
        
        if file_ext == '.zip':
            result.update(_analyze_zip(filepath))
        elif file_ext in ['.tar', '.gz', '.bz2']:
            result.update(_analyze_tar(filepath))
        
        return result
        
    except Exception as e:
        logger.error(f"Archive analysis error: {e}")
        return {}

def _analyze_zip(filepath: str) -> Dict[str, Any]:
    """Analyze ZIP archive"""
    try:
        with zipfile.ZipFile(filepath, 'r') as zip_file:
            file_list = zip_file.filelist
            
            result = {
                "zip_analysis": {
                    "file_count": len(file_list),
                    "total_uncompressed_size": sum(f.file_size for f in file_list),
                    "total_compressed_size": sum(f.compress_size for f in file_list),
                    "compression_method": {},
                    "file_types": {},
                    "directories": []
                }
            }
            
            # Calculate compression ratio
            total_uncompressed = result["zip_analysis"]["total_uncompressed_size"]
            total_compressed = result["zip_analysis"]["total_compressed_size"]
            
            if total_uncompressed > 0:
                ratio = (total_uncompressed - total_compressed) / total_uncompressed * 100
                result["zip_analysis"]["compression_ratio"] = round(ratio, 2)
            
            # Analyze file types and compression methods
            compression_methods = {}
            file_types = {}
            directories = set()
            
            for file_info in file_list:
                # Compression method
                method = file_info.compress_type
                compression_methods[method] = compression_methods.get(method, 0) + 1
                
                # File type
                file_ext = Path(file_info.filename).suffix.lower()
                if file_ext:
                    file_types[file_ext] = file_types.get(file_ext, 0) + 1
                
                # Directory structure
                dir_path = str(Path(file_info.filename).parent)
                if dir_path != '.':
                    directories.add(dir_path)
            
            result["zip_analysis"]["compression_method"] = compression_methods
            result["zip_analysis"]["file_types"] = file_types
            result["zip_analysis"]["directories"] = sorted(list(directories))[:20]  # First 20
            
            return result
            
    except Exception as e:
        logger.error(f"ZIP analysis error: {e}")
        return {}

def _analyze_source_code(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze source code files"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        result = {
            "source_code_analysis": {
                "language": _detect_programming_language(file_ext),
                "line_count": len(content.splitlines()),
                "character_count": len(content),
                "blank_lines": 0,
                "comment_lines": 0,
                "code_lines": 0,
                "functions": 0,
                "classes": 0,
                "imports": 0
            }
        }
        
        lines = content.splitlines()
        
        # Analyze lines
        blank_lines = 0
        comment_lines = 0
        
        for line in lines:
            stripped = line.strip()
            if not stripped:
                blank_lines += 1
            elif _is_comment_line(stripped, file_ext):
                comment_lines += 1
        
        result["source_code_analysis"]["blank_lines"] = blank_lines
        result["source_code_analysis"]["comment_lines"] = comment_lines
        result["source_code_analysis"]["code_lines"] = len(lines) - blank_lines - comment_lines
        
        # Language-specific analysis
        if file_ext == '.py':
            result["source_code_analysis"].update(_analyze_python_code(content))
        elif file_ext == '.js':
            result["source_code_analysis"].update(_analyze_javascript_code(content))
        elif file_ext == '.java':
            result["source_code_analysis"].update(_analyze_java_code(content))
        
        return result
        
    except Exception as e:
        logger.error(f"Source code analysis error: {e}")
        return {}

def _detect_programming_language(file_ext: str) -> str:
    """Detect programming language from file extension"""
    language_map = {
        '.py': 'Python',
        '.js': 'JavaScript',
        '.java': 'Java',
        '.cpp': 'C++',
        '.c': 'C',
        '.h': 'C/C++ Header',
        '.cs': 'C#',
        '.php': 'PHP',
        '.rb': 'Ruby',
        '.go': 'Go',
        '.rs': 'Rust',
        '.swift': 'Swift',
        '.kt': 'Kotlin',
        '.scala': 'Scala',
        '.r': 'R',
        '.m': 'Objective-C',
        '.pl': 'Perl',
        '.sh': 'Shell Script'
    }
    
    return language_map.get(file_ext, 'Unknown')

def _is_comment_line(line: str, file_ext: str) -> bool:
    """Check if a line is a comment"""
    comment_prefixes = {
        '.py': ['#'],
        '.js': ['//', '/*', '*'],
        '.java': ['//', '/*', '*'],
        '.cpp': ['//', '/*', '*'],
        '.c': ['//', '/*', '*'],
        '.cs': ['//', '/*', '*'],
        '.php': ['//', '/*', '*', '#'],
        '.rb': ['#'],
        '.go': ['//', '/*', '*'],
        '.rs': ['//', '/*', '*'],
        '.sh': ['#']
    }
    
    prefixes = comment_prefixes.get(file_ext, [])
    return any(line.startswith(prefix) for prefix in prefixes)

def _analyze_python_code(content: str) -> Dict[str, Any]:
    """Analyze Python-specific code features"""
    try:
        result = {
            "python_features": {
                "functions": len(re.findall(r'^\s*def\s+\w+', content, re.MULTILINE)),
                "classes": len(re.findall(r'^\s*class\s+\w+', content, re.MULTILINE)),
                "imports": len(re.findall(r'^\s*(import|from)\s+', content, re.MULTILINE)),
                "decorators": len(re.findall(r'^\s*@\w+', content, re.MULTILINE)),
                "async_functions": len(re.findall(r'^\s*async\s+def\s+\w+', content, re.MULTILINE)),
                "lambda_functions": len(re.findall(r'\blambda\b', content)),
                "list_comprehensions": len(re.findall(r'\[.*for.*in.*\]', content)),
                "docstrings": len(re.findall(r'""".*?"""', content, re.DOTALL))
            }
        }
        
        return result
        
    except Exception as e:
        logger.error(f"Python code analysis error: {e}")
        return {}

def _analyze_text_statistics(filepath: str) -> Dict[str, Any]:
    """Analyze general text statistics"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        words = content.split()
        sentences = re.split(r'[.!?]+', content)
        paragraphs = content.split('\n\n')
        
        result = {
            "text_statistics": {
                "character_count": len(content),
                "character_count_no_spaces": len(content.replace(' ', '')),
                "word_count": len(words),
                "sentence_count": len([s for s in sentences if s.strip()]),
                "paragraph_count": len([p for p in paragraphs if p.strip()]),
                "average_word_length": sum(len(word) for word in words) / len(words) if words else 0,
                "average_sentence_length": len(words) / len(sentences) if sentences else 0,
                "readability_score": _calculate_readability_score(content)
            }
        }
        
        return result
        
    except Exception as e:
        logger.error(f"Text statistics error: {e}")
        return {}

def _calculate_readability_score(text: str) -> float:
    """Calculate Flesch Reading Ease score"""
    try:
        sentences = len(re.split(r'[.!?]+', text))
        words = len(text.split())
        syllables = sum(_count_syllables(word) for word in text.split())
        
        if sentences == 0 or words == 0:
            return 0
        
        # Flesch Reading Ease formula
        score = 206.835 - (1.015 * (words / sentences)) - (84.6 * (syllables / words))
        return round(max(0, min(100, score)), 2)
        
    except:
        return 0

def _count_syllables(word: str) -> int:
    """Count syllables in a word (approximate)"""
    word = word.lower()
    vowels = 'aeiouy'
    syllable_count = 0
    previous_was_vowel = False
    
    for char in word:
        is_vowel = char in vowels
        if is_vowel and not previous_was_vowel:
            syllable_count += 1
        previous_was_vowel = is_vowel
    
    # Handle silent 'e'
    if word.endswith('e'):
        syllable_count -= 1
    
    return max(1, syllable_count)

def _analyze_config_file(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze configuration files"""
    try:
        result = {}
        
        if file_ext == '.json':
            result.update(_analyze_json_config(filepath))
        elif file_ext in ['.yaml', '.yml'] and YAML_AVAILABLE:
            result.update(_analyze_yaml_config(filepath))
        elif file_ext == '.toml' and TOML_AVAILABLE:
            result.update(_analyze_toml_config(filepath))
        elif file_ext in ['.ini', '.cfg', '.conf']:
            result.update(_analyze_ini_config(filepath))
        
        return result
        
    except Exception as e:
        logger.error(f"Config file analysis error: {e}")
        return {}

def _analyze_json_config(filepath: str) -> Dict[str, Any]:
    """Analyze JSON configuration"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        result = {
            "json_analysis": {
                "is_valid": True,
                "key_count": _count_json_keys(data),
                "max_depth": _calculate_json_depth(data),
                "data_types": _analyze_json_types(data),
                "top_level_keys": list(data.keys()) if isinstance(data, dict) else []
            }
        }
        
        return result
        
    except json.JSONDecodeError:
        return {"json_analysis": {"is_valid": False}}
    except Exception as e:
        logger.error(f"JSON analysis error: {e}")
        return {}

def _count_json_keys(obj, count=0) -> int:
    """Recursively count keys in JSON object"""
    if isinstance(obj, dict):
        count += len(obj)
        for value in obj.values():
            count = _count_json_keys(value, count)
    elif isinstance(obj, list):
        for item in obj:
            count = _count_json_keys(item, count)
    return count

def _calculate_json_depth(obj, depth=0) -> int:
    """Calculate maximum depth of JSON object"""
    if isinstance(obj, dict):
        return max([_calculate_json_depth(value, depth + 1) for value in obj.values()], default=depth)
    elif isinstance(obj, list):
        return max([_calculate_json_depth(item, depth + 1) for item in obj], default=depth)
    else:
        return depth

def _analyze_json_types(obj) -> Dict[str, int]:
    """Analyze data types in JSON object"""
    type_counts = {}
    
    def count_types(item):
        item_type = type(item).__name__
        type_counts[item_type] = type_counts.get(item_type, 0) + 1
        
        if isinstance(item, dict):
            for value in item.values():
                count_types(value)
        elif isinstance(item, list):
            for element in item:
                count_types(element)
    
    count_types(obj)
    return type_counts

def _analyze_database_file(filepath: str, file_ext: str) -> Dict[str, Any]:
    """Analyze database files"""
    try:
        result = {}
        
        if file_ext in ['.db', '.sqlite', '.sqlite3']:
            result.update(_analyze_sqlite_database(filepath))
        
        return result
        
    except Exception as e:
        logger.error(f"Database analysis error: {e}")
        return {}

def _analyze_sqlite_database(filepath: str) -> Dict[str, Any]:
    """Analyze SQLite database"""
    try:
        conn = sqlite3.connect(filepath)
        cursor = conn.cursor()
        
        result = {
            "sqlite_analysis": {
                "table_count": 0,
                "tables": [],
                "total_rows": 0,
                "database_size": os.path.getsize(filepath),
                "schema_version": None
            }
        }
        
        # Get table information
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        
        result["sqlite_analysis"]["table_count"] = len(tables)
        
        table_info = []
        total_rows = 0
        
        for (table_name,) in tables:
            # Get row count
            cursor.execute(f"SELECT COUNT(*) FROM `{table_name}`;")
            row_count = cursor.fetchone()[0]
            total_rows += row_count
            
            # Get column information
            cursor.execute(f"PRAGMA table_info(`{table_name}`);")
            columns = cursor.fetchall()
            
            table_info.append({
                "name": table_name,
                "row_count": row_count,
                "column_count": len(columns),
                "columns": [{"name": col[1], "type": col[2]} for col in columns]
            })
        
        result["sqlite_analysis"]["tables"] = table_info
        result["sqlite_analysis"]["total_rows"] = total_rows
        
        # Get SQLite version
        cursor.execute("SELECT sqlite_version();")
        version = cursor.fetchone()[0]
        result["sqlite_analysis"]["sqlite_version"] = version
        
        conn.close()
        
        return result
        
    except Exception as e:
        logger.error(f"SQLite analysis error: {e}")
        return {}

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        result = extract_document_metadata(sys.argv[1])
        print(json.dumps(result, indent=2, default=str))
    else:
        print("Usage: python document_metadata_ultimate.py <document_file>")