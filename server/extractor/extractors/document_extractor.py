"""
Document metadata extractor for MetaExtract.

Specialized extractor for document file formats including PDF, Office documents,
OpenDocument formats, e-books, web documents, and other document types.
Extracts document properties, content metadata, security information, and
advanced document-specific metadata.
"""

import logging
import json
import traceback
from pathlib import Path
from typing import Any, Dict, Optional, List
from datetime import datetime

from ..core.base_engine import BaseExtractor, ExtractionContext, ExtractionResult, ExtractionStatus

logger = logging.getLogger(__name__)

# Availability flags for optional libraries
try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    try:
        from PyPDF2 import PdfReader
        PYPDF_AVAILABLE = True
    except ImportError:
        PYPDF_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Cm
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


class DocumentExtractor(BaseExtractor):
    """
    Specialized extractor for document file formats.
    
    Supports PDF, Office documents (Word, Excel, PowerPoint), OpenDocument formats,
    e-books, web documents, text files, and other document types. Extracts document
    properties, content metadata, security information, and format-specific details.
    """
    
    def __init__(self):
        """Initialize the document extractor."""
        supported_formats = [
            # PDF Documents
            '.pdf', '.pdfa', '.pdfx',
            
            # Microsoft Office
            '.docx', '.docm', '.doc', '.dotx', '.dotm',
            '.xlsx', '.xlsm', '.xls', '.xltx', '.xltm', '.xlsb',
            '.pptx', '.pptm', '.ppt', '.potx', '.potm', '.ppsx', '.ppsm',
            
            # OpenDocument Format
            '.odt', '.ods', '.odp', '.odg', '.odc', '.odf', '.odi', '.odm',
            
            # E-books
            '.epub', '.mobi', '.azw', '.azw3', '.fb2', '.djvu', '.cbz', '.cbr',
            
            # Web Documents
            '.html', '.htm', '.xhtml', '.xml', '.css', '.js', '.json', '.yaml', '.yml', '.toml',
            
            # Text Documents
            '.txt', '.rtf', '.md', '.markdown', '.rst', '.tex', '.ltx',
            
            # Data Formats
            '.csv', '.tsv', '.psv', '.log',
            
            # Archives (as documents)
            '.zip', '.tar', '.gz', '.bz2', '.xz', '.7z', '.rar',
            
            # Other Document Formats
            '.ps', '.eps', '.ai', '.indd', '.qxp', '.pub', '.vsd', '.vsdx', '.mpp',
            '.one', '.onetoc2', '.notebook'
        ]
        super().__init__("document_extractor", supported_formats)
        self.logger = logging.getLogger(__name__)
    
    def _extract_metadata(self, context: ExtractionContext) -> Dict[str, Any]:
        """
        Extract metadata from a document file.
        
        Args:
            context: Extraction context containing file information
            
        Returns:
            Dictionary containing extracted document metadata
        """
        filepath = context.filepath
        metadata = {}
        
        try:
            # Basic file information
            metadata["file_info"] = self._extract_file_info(filepath)
            
            # Format-specific extraction
            ext = Path(filepath).suffix.lower()
            
            if ext == '.pdf':
                metadata["pdf"] = self._extract_pdf_metadata(filepath)
            elif ext in ['.docx', '.docm', '.doc']:
                metadata["word"] = self._extract_word_metadata(filepath)
            elif ext in ['.xlsx', '.xlsm', '.xls']:
                metadata["excel"] = self._extract_excel_metadata(filepath)
            elif ext in ['.pptx', '.pptm', '.ppt']:
                metadata["powerpoint"] = self._extract_powerpoint_metadata(filepath)
            elif ext in ['.odt', '.ods', '.odp']:
                metadata["opendocument"] = self._extract_opendocument_metadata(filepath)
            elif ext in ['.epub', '.mobi', '.azw', '.azw3']:
                metadata["ebook"] = self._extract_ebook_metadata(filepath)
            elif ext in ['.html', '.htm', '.xhtml']:
                metadata["html"] = self._extract_html_metadata(filepath)
            elif ext in ['.xml', '.json', '.yaml', '.yml']:
                metadata["structured"] = self._extract_structured_data_metadata(filepath)
            elif ext in ['.txt', '.md', '.markdown', '.rst']:
                metadata["text"] = self._extract_text_metadata(filepath)
            elif ext in ['.csv', '.tsv', '.psv']:
                metadata["tabular"] = self._extract_tabular_metadata(filepath)
            else:
                # Generic document metadata
                metadata["generic"] = self._extract_generic_document_metadata(filepath)
            
            # Add extraction statistics
            metadata["extraction_stats"] = {
                "pypdf_available": PYPDF_AVAILABLE,
                "python_docx_available": PYTHON_DOCX_AVAILABLE,
                "openpyxl_available": OPENPYXL_AVAILABLE,
                "python_pptx_available": PYTHON_PPTX_AVAILABLE,
                "format_detected": ext,
                "specialized_extractor_used": self._get_specialized_extractor_used(ext)
            }
            
            return metadata
            
        except Exception as e:
            self.logger.error(f"Document extraction failed for {filepath}: {e}")
            self.logger.debug(f"Traceback: {traceback.format_exc()}")
            raise
    
    def _extract_file_info(self, filepath: str) -> Dict[str, Any]:
        """Extract basic file information."""
        try:
            stat_info = Path(filepath).stat()
            path = Path(filepath)
            
            return {
                "filename": path.name,
                "file_size_bytes": stat_info.st_size,
                "file_extension": path.suffix.lower(),
                "absolute_path": str(path.absolute()),
                "creation_time": stat_info.st_ctime,
                "modification_time": stat_info.st_mtime,
                "access_time": stat_info.st_atime
            }
        except Exception as e:
            self.logger.warning(f"Could not extract file info for {filepath}: {e}")
            return {}
    
    def _extract_pdf_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract PDF-specific metadata."""
        if not PYPDF_AVAILABLE:
            return None
        
        try:
            reader = PdfReader(filepath)
            
            # Basic PDF properties
            pdf_info = {
                "page_count": len(reader.pages),
                "is_encrypted": reader.is_encrypted,
                "metadata_available": bool(reader.metadata)
            }
            
            # PDF metadata
            if reader.metadata:
                metadata = reader.metadata
                pdf_info.update({
                    "title": metadata.get("/Title"),
                    "author": metadata.get("/Author"),
                    "subject": metadata.get("/Subject"),
                    "creator": metadata.get("/Creator"),
                    "producer": metadata.get("/Producer"),
                    "creation_date": str(metadata.get("/CreationDate")),
                    "modification_date": str(metadata.get("/ModDate")),
                    "keywords": metadata.get("/Keywords")
                })
            
            # Page layout information (from first page)
            if reader.pages:
                first_page = reader.pages[0]
                pdf_info["first_page"] = {
                    "width_pts": float(first_page.mediabox.width) if hasattr(first_page, 'mediabox') else None,
                    "height_pts": float(first_page.mediabox.height) if hasattr(first_page, 'mediabox') else None,
                    "rotation": first_page.rotation if hasattr(first_page, 'rotation') else None
                }
            
            return pdf_info
            
        except Exception as e:
            self.logger.warning(f"PDF metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_word_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract Microsoft Word document metadata."""
        if not PYTHON_DOCX_AVAILABLE:
            return None
        
        try:
            doc = Document(filepath)
            
            word_info = {
                "document_type": "Word Document",
                "core_properties_available": hasattr(doc, 'core_properties'),
                "paragraphs_count": len(doc.paragraphs),
                "tables_count": len(doc.tables),
                "sections_count": len(doc.sections)
            }
            
            # Core properties
            if hasattr(doc, 'core_properties'):
                props = doc.core_properties
                word_info.update({
                    "title": props.title,
                    "author": props.author,
                    "subject": props.subject,
                    "keywords": props.keywords,
                    "category": props.category,
                    "comments": props.comments,
                    "created": str(props.created),
                    "modified": str(props.modified),
                    "last_modified_by": props.last_modified_by,
                    "revision": props.revision,
                    "language": props.language,
                    "identifier": props.identifier
                })
            
            # Document structure
            word_info["structure"] = {
                "has_headers": any(section.header for section in doc.sections),
                "has_footers": any(section.footer for section in doc.sections),
                "has_page_numbers": any("page" in str(p.text).lower() for p in doc.paragraphs),
                "has_tables": len(doc.tables) > 0,
                "has_images": len(doc.inline_shapes) > 0
            }
            
            return word_info
            
        except Exception as e:
            self.logger.warning(f"Word metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_excel_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract Microsoft Excel spreadsheet metadata."""
        if not OPENPYXL_AVAILABLE:
            return None
        
        try:
            wb = load_workbook(filepath, read_only=True, data_only=True)
            
            excel_info = {
                "document_type": "Excel Spreadsheet",
                "worksheets_count": len(wb.sheetnames),
                "worksheet_names": wb.sheetnames,
                "has_macros": wb.vba_archive is not None,
                "active_sheet": wb.active.title if wb.active else None
            }
            
            # Worksheet analysis
            worksheets_info = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Basic worksheet properties
                ws_info = {
                    "name": sheet_name,
                    "max_row": ws.max_row,
                    "max_column": ws.max_column,
                    "has_formulas": False,  # Will check below
                    "has_formatting": False,  # Will check below
                    "has_charts": False,  # Will check below
                    "has_pivot_tables": False  # Will check below
                }
                
                # Check for formulas (sample first 10x10 area)
                formula_count = 0
                for row in range(1, min(11, ws.max_row + 1)):
                    for col in range(1, min(11, ws.max_column + 1)):
                        cell = ws.cell(row=row, column=col)
                        if cell.value and str(cell.value).startswith('='):
                            formula_count += 1
                
                ws_info["formula_count"] = formula_count
                ws_info["has_formulas"] = formula_count > 0
                
                worksheets_info.append(ws_info)
            
            excel_info["worksheets"] = worksheets_info
            
            return excel_info
            
        except Exception as e:
            self.logger.warning(f"Excel metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_powerpoint_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract Microsoft PowerPoint presentation metadata."""
        if not PYTHON_PPTX_AVAILABLE:
            return None
        
        try:
            prs = Presentation(filepath)
            
            ppt_info = {
                "document_type": "PowerPoint Presentation",
                "slides_count": len(prs.slides),
                "slide_layouts_available": len(prs.slide_layouts),
                "slide_masters_available": len(prs.slide_masters),
                "has_notes": any(slide.has_notes_slide for slide in prs.slides)
            }
            
            # Slide analysis
            slides_info = []
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "slide_number": i + 1,
                    "has_title": bool(slide.shapes.title),
                    "shapes_count": len(slide.shapes),
                    "has_notes": slide.has_notes_slide,
                    "slide_layout": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown"
                }
                slides_info.append(slide_info)
            
            ppt_info["slides"] = slides_info
            
            return ppt_info
            
        except Exception as e:
            self.logger.warning(f"PowerPoint metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_opendocument_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract OpenDocument format metadata."""
        try:
            # OpenDocument files are ZIP archives with specific structure
            import zipfile
            
            with zipfile.ZipFile(filepath, 'r') as zip_file:
                # Check for OpenDocument structure
                manifest_exists = 'META-INF/manifest.xml' in zip_file.namelist()
                
                od_info = {
                    "format": "OpenDocument",
                    "is_valid_opendocument": manifest_exists,
                    "files_in_archive": len(zip_file.namelist()),
                    "archive_contents": []
                }
                
                # Analyze archive contents
                for file_info in zip_file.filelist:
                    od_info["archive_contents"].append({
                        "filename": file_info.filename,
                        "file_size": file_info.file_size,
                        "compressed_size": file_info.compress_size,
                        "compression_ratio": file_info.file_size / file_info.compress_size if file_info.compress_size > 0 else 0
                    })
                
                # Look for specific OpenDocument files
                content_files = [f for f in zip_file.namelist() if f.startswith('content.xml')]
                styles_files = [f for f in zip_file.namelist() if f.startswith('styles.xml')]
                meta_files = [f for f in zip_file.namelist() if f.startswith('meta.xml')]
                
                od_info["has_content_xml"] = len(content_files) > 0
                od_info["has_styles_xml"] = len(styles_files) > 0
                od_info["has_meta_xml"] = len(meta_files) > 0
                
                return od_info
                
        except (zipfile.BadZipFile, Exception) as e:
            self.logger.warning(f"OpenDocument metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_ebook_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract e-book metadata."""
        try:
            ext = Path(filepath).suffix.lower()
            ebook_info = {
                "format": ext,
                "is_ebook": True
            }
            
            if ext == '.epub':
                # EPUB is a ZIP file with specific structure
                import zipfile
                
                with zipfile.ZipFile(filepath, 'r') as zip_file:
                    # Look for container.xml and OPF file
                    container_exists = 'META-INF/container.xml' in zip_file.namelist()
                    opf_files = [f for f in zip_file.namelist() if f.endswith('.opf')]
                    
                    ebook_info.update({
                        "is_valid_epub": container_exists,
                        "opf_files_count": len(opf_files),
                        "opf_files": opf_files[:3]  # First 3 OPF files
                    })
            
            elif ext == '.mobi':
                # MOBI files have specific header structure
                with open(filepath, 'rb') as f:
                    header = f.read(100)
                    
                    # Check for MOBI signature
                    is_mobi = header.startswith(b'BOOKMOBI') or b'MOBI' in header
                    
                    ebook_info.update({
                        "is_valid_mobi": is_mobi,
                        "file_size": len(header)
                    })
            
            return ebook_info
            
        except Exception as e:
            self.logger.warning(f"E-book metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_html_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract HTML document metadata."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(10000)  # Read first 10KB
            
            html_info = {
                "format": "HTML",
                "file_size_bytes": Path(filepath).stat().st_size,
                "has_doctype": "<!DOCTYPE" in content.upper(),
                "has_html_tag": "<html" in content.lower(),
                "has_head_tag": "<head" in content.lower(),
                "has_body_tag": "<body" in content.lower()
            }
            
            # Extract meta tags
            import re
            meta_tags = re.findall(r'<meta[^>]*>', content, re.IGNORECASE)
            html_info["meta_tags_count"] = len(meta_tags)
            
            # Extract title if present
            title_match = re.search(r'<title[^>]*>([^<]*)</title>', content, re.IGNORECASE)
            html_info["has_title"] = bool(title_match)
            if title_match:
                html_info["title"] = title_match.group(1).strip()
            
            # Extract charset if present
            charset_match = re.search(r'charset=([^"\'>\s]+)', content, re.IGNORECASE)
            html_info["has_charset"] = bool(charset_match)
            if charset_match:
                html_info["charset"] = charset_match.group(1)
            
            return html_info
            
        except Exception as e:
            self.logger.warning(f"HTML metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_structured_data_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract structured data file metadata (XML, JSON, YAML)."""
        try:
            ext = Path(filepath).suffix.lower()
            structured_info = {
                "format": ext,
                "file_size_bytes": Path(filepath).stat().st_size
            }
            
            if ext == '.json':
                import json
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    try:
                        data = json.load(f)
                        structured_info.update({
                            "is_valid_json": True,
                            "json_type": type(data).__name__,
                            "has_nested_structure": isinstance(data, dict) and len(data) > 0
                        })
                    except json.JSONDecodeError:
                        structured_info["is_valid_json"] = False
            
            elif ext == '.xml':
                import xml.etree.ElementTree as ET
                try:
                    tree = ET.parse(filepath)
                    root = tree.getroot()
                    structured_info.update({
                        "is_valid_xml": True,
                        "root_tag": root.tag,
                        "has_attributes": bool(root.attrib),
                        "attributes_count": len(root.attrib)
                    })
                except ET.ParseError:
                    structured_info["is_valid_xml"] = False
            
            elif ext in ['.yaml', '.yml']:
                try:
                    import yaml
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        try:
                            data = yaml.safe_load(f)
                            structured_info.update({
                                "is_valid_yaml": True,
                                "yaml_type": type(data).__name__ if data else "None"
                            })
                        except yaml.YAMLError:
                            structured_info["is_valid_yaml"] = False
                except ImportError:
                    structured_info["yaml_library_available"] = False
            
            return structured_info
            
        except Exception as e:
            self.logger.warning(f"Structured data metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_text_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract text document metadata."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            text_info = {
                "format": "Text Document",
                "file_size_bytes": len(content.encode('utf-8')),
                "character_count": len(content),
                "line_count": content.count('\n') + 1,
                "word_count": len(content.split()),
                "has_content": len(content.strip()) > 0
            }
            
            # Detect encoding issues
            try:
                content.encode('utf-8')
                text_info["encoding_issues"] = False
            except UnicodeEncodeError:
                text_info["encoding_issues"] = True
            
            # Basic content analysis
            text_info["content_analysis"] = {
                "has_numbers": any(c.isdigit() for c in content),
                "has_special_chars": any(not c.isalnum() and not c.isspace() for c in content),
                "has_urls": "http://" in content or "https://" in content,
                "has_email": "@" in content and "." in content
            }
            
            return text_info
            
        except Exception as e:
            self.logger.warning(f"Text metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_tabular_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract tabular data file metadata (CSV, TSV)."""
        try:
            ext = Path(filepath).suffix.lower()
            delimiter = ',' if ext == '.csv' else '\t' if ext == '.tsv' else '|'
            
            tabular_info = {
                "format": ext,
                "delimiter": delimiter,
                "file_size_bytes": Path(filepath).stat().st_size
            }
            
            # Read first few lines to analyze structure
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                lines = []
                for i, line in enumerate(f):
                    if i >= 10:  # Read first 10 lines
                        break
                    lines.append(line.strip())
            
            if lines:
                tabular_info.update({
                    "rows_count": len(lines),
                    "has_header": any(not str(line).replace(delimiter, '').isdigit() for line in lines[:3]),
                    "column_count": max(line.count(delimiter) + 1 for line in lines) if lines else 0,
                    "sample_rows": lines[:3]
                })
            
            return tabular_info
            
        except Exception as e:
            self.logger.warning(f"Tabular metadata extraction failed for {filepath}: {e}")
            return None
    
    def _extract_generic_document_metadata(self, filepath: str) -> Optional[Dict[str, Any]]:
        """Extract generic document metadata for unsupported formats."""
        try:
            # Try to read first few bytes to detect file type
            with open(filepath, 'rb') as f:
                header = f.read(100)
            
            generic_info = {
                "format": "Generic Document",
                "file_size_bytes": Path(filepath).stat().st_size,
                "header_bytes": header[:20].hex(),
                "header_ascii": header[:20].decode('ascii', errors='ignore').replace('\x00', ''),
                "detected_type": self._detect_file_type_from_header(header)
            }
            
            return generic_info
            
        except Exception as e:
            self.logger.warning(f"Generic document metadata extraction failed for {filepath}: {e}")
            return None
    
    def _detect_file_type_from_header(self, header: bytes) -> str:
        """Detect file type from header bytes."""
        # Magic numbers for common file types
        if header.startswith(b'%PDF'):
            return 'PDF'
        elif header.startswith(b'PK'):
            return 'ZIP/Office Document'
        elif header.startswith(b'\xff\xd8'):
            return 'JPEG Image'
        elif header.startswith(b'\x89PNG'):
            return 'PNG Image'
        elif header.startswith(b'GIF'):
            return 'GIF Image'
        elif header.startswith(b'BM'):
            return 'BMP Image'
        elif header.startswith(b'RIFF') and b'WEBP' in header:
            return 'WebP Image'
        elif header.startswith(b'\x00\x00\x00\x14ftyp'):
            return 'MP4 Video'
        elif header.startswith(b'ID3'):
            return 'MP3 Audio'
        elif header.startswith(b'fLaC'):
            return 'FLAC Audio'
        elif header.startswith(b'OggS'):
            return 'OGG Audio'
        elif header.startswith(b'\x1a\x45\xdf\xa3'):
            return 'WebM Video'
        elif header.startswith(b'{') or header.startswith(b'['):
            return 'JSON/Text'
        elif header.startswith(b'<'):
            return 'XML/HTML'
        else:
            return 'Unknown'
    
    def _get_specialized_extractor_used(self, ext: str) -> str:
        """Get the name of the specialized extractor used."""
        extractors = {
            '.pdf': 'pdf',
            '.docx': 'word', '.docm': 'word', '.doc': 'word',
            '.xlsx': 'excel', '.xlsm': 'excel', '.xls': 'excel',
            '.pptx': 'powerpoint', '.pptm': 'powerpoint', '.ppt': 'powerpoint',
            '.odt': 'opendocument', '.ods': 'opendocument', '.odp': 'opendocument',
            '.epub': 'ebook', '.mobi': 'ebook', '.azw': 'ebook', '.azw3': 'ebook',
            '.html': 'html', '.htm': 'html', '.xhtml': 'html',
            '.xml': 'structured', '.json': 'structured', '.yaml': 'structured', '.yml': 'structured',
            '.txt': 'text', '.md': 'text', '.markdown': 'text', '.rst': 'text',
            '.csv': 'tabular', '.tsv': 'tabular', '.psv': 'tabular'
        }
        
        return extractors.get(ext, 'generic')
    
    def get_extraction_info(self) -> Dict[str, Any]:
        """Get information about this extractor."""
        return {
            "name": self.name,
            "supported_formats": self.supported_formats,
            "capabilities": {
                "pypdf": PYPDF_AVAILABLE,
                "python_docx": PYTHON_DOCX_AVAILABLE,
                "openpyxl": OPENPYXL_AVAILABLE,
                "python_pptx": PYTHON_PPTX_AVAILABLE,
                "pdf_extraction": PYPDF_AVAILABLE,
                "office_extraction": PYTHON_DOCX_AVAILABLE or OPENPYXL_AVAILABLE or PYTHON_PPTX_AVAILABLE,
                "opendocument_extraction": True,
                "ebook_extraction": True,
                "html_extraction": True,
                "structured_data_extraction": True,
                "text_extraction": True,
                "tabular_extraction": True
            },
            "version": "1.0.0"
        }